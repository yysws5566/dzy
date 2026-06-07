#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
紫金矿业 + 洛阳钼业 箱体波动策略 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 全局样式 ────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # 深蓝黑
CARD_BG   = RGBColor(0x25, 0x25, 0x3D)   # 卡片底色
ACCENT_RED    = RGBColor(0xFF, 0x45, 0x45)   # 卖出/阻力
ACCENT_GREEN  = RGBColor(0x00, 0xD2, 0x8A)   # 买入/支撑
ACCENT_GOLD   = RGBColor(0xFF, 0xD7, 0x00)   # 高亮
ACCENT_BLUE   = RGBColor(0x4D, 0xAB, 0xF7)   # 信息
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
TITLE_FONT = "Microsoft YaHei"
BODY_FONT  = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

# ── 工具函数 ────────────────────────────────────────────
def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_card(slide, left, top, width, height, color=CARD_BG, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, font_name=BODY_FONT, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=13, color=WHITE, line_spacing=1.3):
    """lines: list of (text, bold, color_override) or str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, bold, clr = line, False, color
        else:
            text, bold, clr = line[0], line[1] if len(line) > 1 else False, line[2] if len(line) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = BODY_FONT
        p.space_after = Pt(font_size * (line_spacing - 1) * 0.5)
    return txBox

def add_section_title(slide, left, top, title, subtitle=None):
    """模块标题：大号标题 + 可选副标题"""
    add_text_box(slide, left, top, Inches(10), Inches(0.6), title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, left, top + Inches(0.55), Inches(10), Inches(0.4), subtitle, font_size=14, color=LIGHT_GRAY)

def add_tag(slide, left, top, text, bg_color, text_color=WHITE, font_size=11):
    """小标签"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.4), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = BODY_FONT
    return shape

# ══════════════════════════════════════════════════════════
# 第1页：封面
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
# 装饰条
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_GOLD; bar.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(0.8),
             "📦 箱体波动高胜率策略", font_size=42, bold=True, color=WHITE)
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.5),
             "紫金矿业（601899）  ×  洛阳钼业（603993）", font_size=20, color=LIGHT_GRAY)
add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.5),
             "模块化操作手册  |  支撑/阻力定量  |  日内/周线买卖点  |  仓位管理", font_size=16, color=ACCENT_BLUE)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4),
             "数据截至 2026年6月5日收盘  ·  仅供思路参考，不构成投资建议", font_size=12, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# 第2页：模块总览（目录）
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, Inches(0.8), Inches(0.5), "📋 策略模块总览", "点击跳转至对应模块")

modules = [
    ("M1", "市场定位", "两只票当前状态速览", ACCENT_BLUE),
    ("M2", "紫金矿业箱体", "29-31.5元核心箱体 + 三区买卖点", ACCENT_GREEN),
    ("M3", "紫金买卖规则", "分档买入/卖出 + 日内周线阈值", ACCENT_GOLD),
    ("M4", "洛阳钼业箱体", "18-19.5元核心箱体 + 三区买卖点", ACCENT_GREEN),
    ("M5", "洛钼买卖规则", "分档买入/卖出 + 日内周线阈值", ACCENT_GOLD),
    ("M6", "风控体系", "止损/仓位/成交量确认/大盘配合", ACCENT_RED),
    ("M7", "当前作战地图", "两只票现在处于什么位置、该做什么", WHITE),
]

for i, (code, title, desc, color) in enumerate(modules):
    y = Inches(1.6) + Inches(0.75) * i
    add_card(slide, Inches(1.2), y, Inches(10.8), Inches(0.62), color=CARD_BG, border=color)
    add_text_box(slide, Inches(1.5), y + Inches(0.08), Inches(1.0), Inches(0.45),
                 code, font_size=18, bold=True, color=color)
    add_text_box(slide, Inches(2.6), y + Inches(0.08), Inches(3.5), Inches(0.45),
                 title, font_size=16, bold=True, color=WHITE)
    add_text_box(slide, Inches(6.2), y + Inches(0.10), Inches(5.5), Inches(0.45),
                 desc, font_size=12, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# 第3页：M1 - 市场定位速览
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, Inches(0.8), Inches(0.3), "M1  市场定位", "两只股票当前状态速览 — 2026.06.05 收盘")

# 紫金卡片
card = add_card(slide, Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.5), color=CARD_BG, border=ACCENT_GOLD)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
             "🏔️ 紫金矿业  601899", font_size=22, bold=True, color=ACCENT_GOLD)
add_multiline_box(slide, Inches(0.8), Inches(2.1), Inches(5.3), Inches(4.5), [
    ("📌 最新价：29.63 元  |  日内振幅 4.35%", True, WHITE),
    ("",),
    ("📊 年化波动率：35.36%（偏高波动，适合做T）", False, ACCENT_BLUE),
    ("📊 动态PE：9.81 倍（低估值 + 基本面扎实）", False, ACCENT_BLUE),
    ("",),
    ("🔴 技术面：MACD零轴下死叉，偏空", False, ACCENT_RED),
    ("🔴 资金面：主力10日净流出62.92亿", False, ACCENT_RED),
    ("🔴 筹码面：88%持仓者亏损，均价33.01元", False, ACCENT_RED),
    ("",),
    ("💡 定位：箱体下沿附近，等止跌确认信号", False, ACCENT_GREEN),
    ("   核心箱体 → 29.00 ~ 31.50 元", False, WHITE),
], font_size=13)

# 洛钼卡片
card = add_card(slide, Inches(6.9), Inches(1.3), Inches(5.9), Inches(5.5), color=CARD_BG, border=ACCENT_GOLD)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "⛏️ 洛阳钼业  603993", font_size=22, bold=True, color=ACCENT_GOLD)
add_multiline_box(slide, Inches(7.2), Inches(2.1), Inches(5.3), Inches(4.5), [
    ("📌 最新价：18.19 元  |  日内振幅 ~5.5%", True, WHITE),
    ("",),
    ("📊 动态PE：12.54 倍", False, ACCENT_BLUE),
    ("📊 Q1净利润：77.6亿（同比+96.65%）", False, ACCENT_BLUE),
    ("",),
    ("🔴 技术面：MACD零轴下死叉，一阴穿5线", False, ACCENT_RED),
    ("🔴 资金面：主力净流出约6亿", False, ACCENT_RED),
    ("🔴 筹码面：套牢盘集中在18-19.5元", False, ACCENT_RED),
    ("",),
    ("💡 定位：箱体下沿附近，同样等止跌信号", False, ACCENT_GREEN),
    ("   核心箱体 → 18.00 ~ 19.50 元", False, WHITE),
], font_size=13)

# ══════════════════════════════════════════════════════════
# 第4页：M2 - 紫金矿业箱体
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, Inches(0.8), Inches(0.3), "M2  紫金矿业 · 箱体三区划分", "核心箱体 29.00~31.50 | 扩展箱体 27.80~34.30")

# 箱体图 — 用色块模拟
levels = [
    ("34.30  🔴 强阻力", 0.55, ACCENT_RED, "突破则趋势反转，空头策略全部失效"),
    ("33.00  🟠 筹码成本线（S3清仓）", 0.6, RGBColor(0xFF, 0x8C, 0x00), "88%持仓者成本位，抛压最重"),
    ("32.00  🟡 箱体上沿（S2卖出）", 0.65, ACCENT_GOLD, "近期反弹的目标上限"),
    ("31.50  🟢 波动上沿（S1减仓）", 0.60, ACCENT_GREEN, "日内涨幅≥2%+到此→减仓1/3"),
    ("──────────────── 核心波动区间 ────────────────", 0.03, LIGHT_GRAY, ""),
    ("29.60  🟢 箱体下沿（B1试探）← 当前价", 0.70, ACCENT_GREEN, "日内跌幅≥2%+到此→试探性买入20%"),
    ("29.00  🟡 整数关口（B2加仓）", 0.65, ACCENT_GOLD, "整数强支撑，胜率最高，加仓30%"),
    ("27.80  🟠 K线支撑（B3重仓）", 0.60, RGBColor(0xFF, 0x8C, 0x00), "放量止跌→重仓50%，盈亏比最优"),
    ("27.50  🔵 止损线", 0.45, ACCENT_BLUE, "收盘跌破→无条件全清，不恋战"),
]

for i, (label, h_factor, color, desc) in enumerate(levels):
    y = Inches(1.3) + Inches(0.52) * i
    w = Inches(8.5) * h_factor
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y, w, Inches(0.42))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    bar.text_frame.word_wrap = True
    bar.text_frame.paragraphs[0].text = label
    bar.text_frame.paragraphs[0].font.size = Pt(13)
    bar.text_frame.paragraphs[0].font.bold = True
    bar.text_frame.paragraphs[0].font.color.rgb = WHITE if color != ACCENT_GOLD else DARK_BG
    bar.text_frame.paragraphs[0].font.name = BODY_FONT
    if desc:
        add_text_box(slide, Inches(1.0) + w + Inches(0.2), y + Inches(0.03), Inches(4.2), Inches(0.38),
                     desc, font_size=11, color=LIGHT_GRAY)

# 右侧波动率卡片
card = add_card(slide, Inches(9.8), Inches(1.3), Inches(3.2), Inches(2.0), color=CARD_BG, border=ACCENT_BLUE)
add_multiline_box(slide, Inches(10.0), Inches(1.5), Inches(2.8), Inches(1.6), [
    ("📈 波动数据", True, ACCENT_BLUE),
    ("",),
    ("年化波动率：35.36%", False, WHITE),
    ("日均振幅：~4%", False, WHITE),
    ("",),
    ("日内涨跌分布：", False, LIGHT_GRAY),
    ("  0-1%：36.9%  |  1-2%：27.8%", False, WHITE),
    ("  2-3%：14.3%  |  3-5%：12.7%", False, WHITE),
    ("  5%+：5.9%", False, WHITE),
], font_size=11)

# 底部总结
add_text_box(slide, Inches(1.0), Inches(6.5), Inches(11), Inches(0.5),
             "💡 核心逻辑：29以下分批买 → 31以上分批卖，箱体不破就反复做。跌破27.5认输。",
             font_size=14, bold=True, color=ACCENT_GOLD)

# ══════════════════════════════════════════════════════════
# 第5页：M3 - 紫金买卖规则
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, Inches(0.8), Inches(0.3), "M3  紫金矿业 · 买卖规则", "分档进场/离场 + 日内周线阈值 + 仓位管理")

# 买入规则表
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.4), "🟢 买入规则", font_size=20, bold=True, color=ACCENT_GREEN)

buy_rules = [
    ("B1 试探买", "29.60元 + 日内跌≥2%", "20%仓位", "箱体下沿第一次接，轻仓试探"),
    ("B2 标准买", "29.00元 + 分时企稳", "30%仓位", "整数关口，胜率最高的买点"),
    ("B3 重仓买", "27.80-28.00 + 放量止跌", "50%仓位", "恐慌低点，盈亏比最优"),
]
for i, (level, trigger, position, note) in enumerate(buy_rules):
    y = Inches(1.8) + Inches(0.7) * i
    card = add_card(slide, Inches(0.8), y, Inches(5.5), Inches(0.58), color=CARD_BG, border=ACCENT_GREEN)
    add_text_box(slide, Inches(1.0), y + Inches(0.10), Inches(1.3), Inches(0.35), level, font_size=14, bold=True, color=ACCENT_GREEN)
    add_text_box(slide, Inches(2.3), y + Inches(0.10), Inches(2.5), Inches(0.35), trigger, font_size=12, color=WHITE)
    add_text_box(slide, Inches(4.8), y + Inches(0.10), Inches(1.2), Inches(0.35), position, font_size=12, bold=True, color=ACCENT_GOLD)
    add_text_box(slide, Inches(1.0), y + Inches(0.35), Inches(5.3), Inches(0.22), note, font_size=10, color=LIGHT_GRAY)

# 卖出规则表
add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5), Inches(0.4), "🔴 卖出规则", font_size=20, bold=True, color=ACCENT_RED)

sell_rules = [
    ("S1 减仓卖", "31.50元 + 日内涨≥2%", "卖出1/3", "箱体上沿先锁利"),
    ("S2 标准卖", "32.00-32.50元", "再卖1/3", "临近密集套牢区"),
    ("S3 清仓卖", "33.00元附近", "全部清仓", "筹码成本线，抛压最大"),
]
for i, (level, trigger, position, note) in enumerate(sell_rules):
    y = Inches(1.8) + Inches(0.7) * i
    card = add_card(slide, Inches(7.0), y, Inches(5.5), Inches(0.58), color=CARD_BG, border=ACCENT_RED)
    add_text_box(slide, Inches(7.2), y + Inches(0.10), Inches(1.3), Inches(0.35), level, font_size=14, bold=True, color=ACCENT_RED)
    add_text_box(slide, Inches(8.5), y + Inches(0.10), Inches(2.5), Inches(0.35), trigger, font_size=12, color=WHITE)
    add_text_box(slide, Inches(11.0), y + Inches(0.10), Inches(1.2), Inches(0.35), position, font_size=12, bold=True, color=ACCENT_GOLD)
    add_text_box(slide, Inches(7.2), y + Inches(0.35), Inches(5.3), Inches(0.22), note, font_size=10, color=LIGHT_GRAY)

# 涨跌幅阈值卡片
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4), "📅 日/周涨跌幅阈值", font_size=18, bold=True, color=ACCENT_BLUE)

threshold_data = [
    ("📉 日内跌超 3% + 触碰箱体下沿(≤29.00)", "→ 买点确认", ACCENT_GREEN),
    ("📈 日内涨超 3% + 触碰箱体上沿(≥31.50)", "→ 卖点确认", ACCENT_RED),
    ("📉 周跌幅 5-7% + 到达箱底(28-29区间)", "→ 中线买点", ACCENT_GREEN),
    ("📈 周涨幅 5-7% + 到达箱顶(31-33区间)", "→ 中线卖点", ACCENT_RED),
]
for i, (condition, action, color) in enumerate(threshold_data):
    y = Inches(4.8) + Inches(0.45) * i
    add_text_box(slide, Inches(1.0), y, Inches(7), Inches(0.35), condition, font_size=12, color=WHITE)
    add_text_box(slide, Inches(8.2), y, Inches(4), Inches(0.35), action, font_size=12, bold=True, color=color)

# ⚠️ 止损
add_card(slide, Inches(0.8), Inches(6.8), Inches(11.8), Inches(0.42), color=RGBColor(0x3D, 0x1A, 0x1A), border=ACCENT_RED)
add_text_box(slide, Inches(1.0), Inches(6.85), Inches(11.5), Inches(0.35),
             "⛔ 止损铁律：收盘价跌破 27.50 元 → 无条件全部清仓，不犹豫、不补仓、不幻想",
             font_size=13, bold=True, color=ACCENT_RED)

# ══════════════════════════════════════════════════════════
# 第6页：M4 - 洛阳钼业箱体
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, Inches(0.8), Inches(0.3), "M4  洛阳钼业 · 箱体三区划分", "核心箱体 18.00~19.50 | 扩展箱体 17.00~20.00")

levels2 = [
    ("20.00  🔴 心理大关（S3清仓）", 0.55, ACCENT_RED, "前期冲关失败点，心理阻力极强"),
    ("19.50  🟠 箱体上沿（S2卖出）", 0.60, RGBColor(0xFF, 0x8C, 0x00), "筹码密集区上沿，解套抛压重"),
    ("19.00  🟡 均线反压区（S1减仓）", 0.65, ACCENT_GOLD, "5/10/20/30日均线密集，反压明显"),
    ("──────────────── 核心波动区间 ────────────────", 0.03, LIGHT_GRAY, ""),
    ("18.20  🟢 箱体下沿（B1试探）← 当前价附近", 0.70, ACCENT_GREEN, "日内跌≥2%+到此→试探买20%"),
    ("18.00  🟡 整数关口（B2加仓）", 0.65, ACCENT_GOLD, "缩量止跌确认→加仓30%"),
    ("17.00  🟠 强支撑（B3重仓）", 0.60, RGBColor(0xFF, 0x8C, 0x00), "放量长下影→重仓50%"),
    ("15.80  🔵 止损线", 0.45, ACCENT_BLUE, "收盘跌破→无条件全清"),
]

for i, (label, h_factor, color, desc) in enumerate(levels2):
    y = Inches(1.3) + Inches(0.52) * i
    w = Inches(8.5) * h_factor
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y, w, Inches(0.42))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    bar.text_frame.word_wrap = True
    bar.text_frame.paragraphs[0].text = label
    bar.text_frame.paragraphs[0].font.size = Pt(13)
    bar.text_frame.paragraphs[0].font.bold = True
    bar.text_frame.paragraphs[0].font.color.rgb = WHITE if color not in [ACCENT_GOLD, RGBColor(0xFF, 0x8C, 0x00)] else DARK_BG
    bar.text_frame.paragraphs[0].font.name = BODY_FONT
    if desc:
        add_text_box(slide, Inches(1.0) + w + Inches(0.2), y + Inches(0.03), Inches(4.2), Inches(0.38),
                     desc, font_size=11, color=LIGHT_GRAY)

# 右侧
card = add_card(slide, Inches(9.8), Inches(1.3), Inches(3.2), Inches(2.0), color=CARD_BG, border=ACCENT_BLUE)
add_multiline_box(slide, Inches(10.0), Inches(1.5), Inches(2.8), Inches(1.6), [
    ("📈 波动数据", True, ACCENT_BLUE),
    ("",),
    ("日均振幅：~5.5%", False, WHITE),
    ("（振幅大于紫金，做T空间更大）", False, LIGHT_GRAY),
    ("",),
    ("Q1业绩：+96.65% YoY", False, ACCENT_GREEN),
    ("动态PE：12.54倍", False, WHITE),
    ("总市值：3892亿", False, WHITE),
    ("",),
    ("⚠️ 6月5日主力净流出约6亿", False, ACCENT_RED),
], font_size=11)

add_text_box(slide, Inches(1.0), Inches(6.5), Inches(11), Inches(0.5),
             "💡 核心逻辑：18以下分批买 → 19以上分批卖，箱体不破就反复做。跌破15.8认输。",
             font_size=14, bold=True, color=ACCENT_GOLD)

# ══════════════════════════════════════════════════════════
# 第7页：M5 - 洛钼买卖规则
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, Inches(0.8), Inches(0.3), "M5  洛阳钼业 · 买卖规则", "分档进场/离场 + 日内周线阈值 + 仓位管理")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.4), "🟢 买入规则", font_size=20, bold=True, color=ACCENT_GREEN)
buy_rules2 = [
    ("B1 试探买", "18.20元 + 日内跌≥2%", "20%仓位", "箱体下沿第一次接，轻仓试探"),
    ("B2 标准买", "18.00元 + 缩量止跌", "30%仓位", "整数关口，支撑最稳的买点"),
    ("B3 重仓买", "17.00-17.50 + 放量长下影", "50%仓位", "恐慌低吸，盈亏比最优"),
]
for i, (level, trigger, position, note) in enumerate(buy_rules2):
    y = Inches(1.8) + Inches(0.7) * i
    card = add_card(slide, Inches(0.8), y, Inches(5.5), Inches(0.58), color=CARD_BG, border=ACCENT_GREEN)
    add_text_box(slide, Inches(1.0), y + Inches(0.10), Inches(1.3), Inches(0.35), level, font_size=14, bold=True, color=ACCENT_GREEN)
    add_text_box(slide, Inches(2.3), y + Inches(0.10), Inches(2.5), Inches(0.35), trigger, font_size=12, color=WHITE)
    add_text_box(slide, Inches(4.8), y + Inches(0.10), Inches(1.2), Inches(0.35), position, font_size=12, bold=True, color=ACCENT_GOLD)
    add_text_box(slide, Inches(1.0), y + Inches(0.35), Inches(5.3), Inches(0.22), note, font_size=10, color=LIGHT_GRAY)

add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5), Inches(0.4), "🔴 卖出规则", font_size=20, bold=True, color=ACCENT_RED)
sell_rules2 = [
    ("S1 减仓卖", "19.00元 + 日内涨≥2%", "卖出1/3", "均线反压区，先锁利"),
    ("S2 标准卖", "19.50元", "再卖1/3", "箱体上沿 + 筹码密集区"),
    ("S3 清仓卖", "20.00元附近", "全部清仓", "心理关口 + 前期冲关失败点"),
]
for i, (level, trigger, position, note) in enumerate(sell_rules2):
    y = Inches(1.8) + Inches(0.7) * i
    card = add_card(slide, Inches(7.0), y, Inches(5.5), Inches(0.58), color=CARD_BG, border=ACCENT_RED)
    add_text_box(slide, Inches(7.2), y + Inches(0.10), Inches(1.3), Inches(0.35), level, font_size=14, bold=True, color=ACCENT_RED)
    add_text_box(slide, Inches(8.5), y + Inches(0.10), Inches(2.5), Inches(0.35), trigger, font_size=12, color=WHITE)
    add_text_box(slide, Inches(11.0), y + Inches(0.10), Inches(1.2), Inches(0.35), position, font_size=12, bold=True, color=ACCENT_GOLD)
    add_text_box(slide, Inches(7.2), y + Inches(0.35), Inches(5.3), Inches(0.22), note, font_size=10, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4), "📅 日/周涨跌幅阈值", font_size=18, bold=True, color=ACCENT_BLUE)
threshold_data2 = [
    ("📉 日内跌超 3% + 触碰箱体下沿(≤18.00)", "→ 买点确认", ACCENT_GREEN),
    ("📈 日内涨超 3% + 触碰箱体上沿(≥19.00)", "→ 卖点确认", ACCENT_RED),
    ("📉 周跌幅 5-8% + 到达箱底(17-18区间)", "→ 中线买点", ACCENT_GREEN),
    ("📈 周涨幅 5-8% + 到达箱顶(19-20区间)", "→ 中线卖点", ACCENT_RED),
]
for i, (condition, action, color) in enumerate(threshold_data2):
    y = Inches(4.8) + Inches(0.45) * i
    add_text_box(slide, Inches(1.0), y, Inches(7), Inches(0.35), condition, font_size=12, color=WHITE)
    add_text_box(slide, Inches(8.2), y, Inches(4), Inches(0.35), action, font_size=12, bold=True, color=color)

add_card(slide, Inches(0.8), Inches(6.8), Inches(11.8), Inches(0.42), color=RGBColor(0x3D, 0x1A, 0x1A), border=ACCENT_RED)
add_text_box(slide, Inches(1.0), Inches(6.85), Inches(11.5), Inches(0.35),
             "⛔ 止损铁律：收盘价跌破 15.80 元 → 无条件全部清仓，不犹豫、不补仓、不幻想",
             font_size=13, bold=True, color=ACCENT_RED)

# ══════════════════════════════════════════════════════════
# 第8页：M6 - 风控体系
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, Inches(0.8), Inches(0.3), "M6  风控体系", "成交量确认 + 大盘配合 + 仓位纪律 + 时间周期")

# 四大模块
# 1. 成交量确认
card = add_card(slide, Inches(0.5), Inches(1.3), Inches(5.9), Inches(2.5), color=CARD_BG, border=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4), "📊 成交量确认（关键！）", font_size=18, bold=True, color=ACCENT_BLUE)
add_multiline_box(slide, Inches(0.8), Inches(2.0), Inches(5.3), Inches(1.6), [
    ("🟢 箱底买入 → 缩量止跌 或 放量长下影", False, ACCENT_GREEN),
    ("   缩量 = 抛压衰竭，放量长下影 = 有资金抄底", False, LIGHT_GRAY),
    ("",),
    ("🔴 箱顶卖出 → 放量滞涨 或 冲高回落", False, ACCENT_RED),
    ("   放量却涨不动 = 主力在出货", False, LIGHT_GRAY),
    ("",),
    ("⚠️ 假突破过滤 → 突破箱体但无量 = 假突破，不追", False, ACCENT_GOLD),
], font_size=12)

# 2. 大盘配合
card = add_card(slide, Inches(6.9), Inches(1.3), Inches(5.9), Inches(2.5), color=CARD_BG, border=ACCENT_BLUE)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.4), "🌍 大盘配合", font_size=18, bold=True, color=ACCENT_BLUE)
add_multiline_box(slide, Inches(7.2), Inches(2.0), Inches(5.3), Inches(1.6), [
    ("🟢 大盘跌 + 个股跌到箱底 → 买点质量更高", False, ACCENT_GREEN),
    ("   系统性恐慌杀跌，往往是黄金坑", False, LIGHT_GRAY),
    ("",),
    ("🔴 大盘涨 + 个股涨到箱顶 → 卖点质量更高", False, ACCENT_RED),
    ("   跟风冲顶，往往是最佳止盈时机", False, LIGHT_GRAY),
    ("",),
    ("⚠️ 大盘涨 + 个股不跟涨 → 警惕弱势，不宜加仓", False, ACCENT_GOLD),
], font_size=12)

# 3. 仓位纪律
card = add_card(slide, Inches(0.5), Inches(4.1), Inches(5.9), Inches(2.5), color=CARD_BG, border=ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4), "💰 仓位纪律", font_size=18, bold=True, color=ACCENT_GREEN)
add_multiline_box(slide, Inches(0.8), Inches(4.8), Inches(5.3), Inches(1.6), [
    ("📌 单票上限：不超过总资金的 50%", False, WHITE),
    ("📌 单次买入：分3批，绝不一次打完", False, WHITE),
    ("📌 止盈后：至少等3天 或 回到箱底再进", False, WHITE),
    ("📌 浮亏不加仓（除非打到B2/B3触发价）", False, WHITE),
    ("📌 连续止损2次 → 暂停该票操作1周，重新评估", False, ACCENT_RED),
], font_size=12)

# 4. 时间周期
card = add_card(slide, Inches(6.9), Inches(4.1), Inches(5.9), Inches(2.5), color=CARD_BG, border=ACCENT_GREEN)
add_text_box(slide, Inches(7.2), Inches(4.2), Inches(5), Inches(0.4), "⏱️ 时间周期参考", font_size=18, bold=True, color=ACCENT_GREEN)
add_multiline_box(slide, Inches(7.2), Inches(4.8), Inches(5.3), Inches(1.6), [
    ("📌 日内做T：振幅4-5%，有空间但需盯盘", False, WHITE),
    ("📌 3-5日波段（⭐推荐）：箱体内来回做", False, ACCENT_GOLD),
    ("    周期短、频率高、胜率稳", False, LIGHT_GRAY),
    ("📌 周线级别：周跌5-8%买 → 周涨5-8%卖", False, WHITE),
    ("    中线持仓，1-3周一个来回", False, LIGHT_GRAY),
    ("📌 月线级别：不做，等箱体突破/跌破再重新判断", False, WHITE),
], font_size=12)

# ══════════════════════════════════════════════════════════
# 第9页：M7 - 当前作战地图
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, Inches(0.8), Inches(0.3), "M7  当前作战地图", "2026年6月5日收盘 — 现在该做什么？")

# 紫金
card = add_card(slide, Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.2), color=CARD_BG, border=ACCENT_GOLD)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.45),
             "🏔️ 紫金矿业 · 当前位置", font_size=20, bold=True, color=ACCENT_GOLD)

# 当前位置指示器
pos_card = add_card(slide, Inches(0.8), Inches(2.1), Inches(5.3), Inches(1.0), color=RGBColor(0x2D, 0x3D, 0x2D), border=ACCENT_GREEN)
add_multiline_box(slide, Inches(1.1), Inches(2.2), Inches(4.8), Inches(0.8), [
    ("📍 当前价 29.63 元  →  处于 箱体下沿（29.60）附近", True, ACCENT_GREEN),
    ("    距离 B1试探买点（29.60）仅 0.03 元", False, WHITE),
    ("    距离 B2加仓点（29.00）还差 0.63 元 ≈ -2.1%", False, LIGHT_GRAY),
    ("    距离 止损线（27.50）还差 2.13 元 ≈ -7.2%", False, ACCENT_RED),
], font_size=12)

add_multiline_box(slide, Inches(0.8), Inches(3.3), Inches(5.3), Inches(3.0), [
    ("✅ 当前建议：观察区", True, ACCENT_BLUE),
    ("",),
    ("1. 已在B1附近，可以试探性买入（20%仓位）", False, WHITE),
    ("   但需要日内跌幅≥2%或分时止跌确认", False, LIGHT_GRAY),
    ("",),
    ("2. 如果不急，等 29.00 再动手（B2标准买）", False, WHITE),
    ("   盈亏比更好，胜率更高", False, LIGHT_GRAY),
    ("",),
    ("3. 止跌确认信号：", False, ACCENT_GOLD),
    ("   ▸ 缩量十字星（抛压衰竭）", False, WHITE),
    ("   ▸ 长下影线（有资金抄底）", False, WHITE),
    ("   ▸ 放量阳线反包前一日阴线", False, WHITE),
    ("",),
    ("4. 如果继续跌到27.80 → 大胆重仓（B3）", False, ACCENT_GREEN),
    ("5. 如果跌破27.50收盘 → 放弃做多，等新箱体", False, ACCENT_RED),
], font_size=12)

# 洛钼
card = add_card(slide, Inches(6.9), Inches(1.3), Inches(5.9), Inches(5.2), color=CARD_BG, border=ACCENT_GOLD)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.45),
             "⛏️ 洛阳钼业 · 当前位置", font_size=20, bold=True, color=ACCENT_GOLD)

pos_card = add_card(slide, Inches(7.2), Inches(2.1), Inches(5.3), Inches(1.0), color=RGBColor(0x2D, 0x3D, 0x2D), border=ACCENT_GREEN)
add_multiline_box(slide, Inches(7.5), Inches(2.2), Inches(4.8), Inches(0.8), [
    ("📍 当前价 18.19 元  →  处于 箱体下沿（18.20）附近", True, ACCENT_GREEN),
    ("    距离 B1试探买点（18.20）仅 0.01 元", False, WHITE),
    ("    距离 B2加仓点（18.00）还差 0.19 元 ≈ -1.0%", False, LIGHT_GRAY),
    ("    距离 止损线（15.80）还差 2.39 元 ≈ -13.1%", False, ACCENT_RED),
], font_size=12)

add_multiline_box(slide, Inches(7.2), Inches(3.3), Inches(5.3), Inches(3.0), [
    ("✅ 当前建议：观察区", True, ACCENT_BLUE),
    ("",),
    ("1. 已在B1附近，可试探性买入（20%仓位）", False, WHITE),
    ("   但同样需要止跌确认信号", False, LIGHT_GRAY),
    ("",),
    ("2. 如果到18.00 → 加仓到50%（B1+B2）", False, WHITE),
    ("   整数关口支撑通常很强", False, LIGHT_GRAY),
    ("",),
    ("3. 止跌确认信号同上：", False, ACCENT_GOLD),
    ("   ▸ 缩量十字星 / 长下影 / 放量反包", False, WHITE),
    ("",),
    ("4. 如果跌到17.00 → 重仓（B3）", False, ACCENT_GREEN),
    ("5. 如果跌破15.80 → 止损，放弃", False, ACCENT_RED),
    ("",),
    ("⚠️ 洛钼振幅5.5% > 紫金4%，做T空间更大", False, ACCENT_BLUE),
    ("   但也意味着波动风险更高，仓位要更谨慎", False, LIGHT_GRAY),
], font_size=12)

# ══════════════════════════════════════════════════════════
# 第10页：总结 + 操作清单
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
# 装饰条
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.9), Inches(13.333), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_GOLD; bar.line.fill.background()

add_section_title(slide, Inches(0.8), Inches(0.3), "🎯 一句话策略总结")

add_multiline_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.2), [
    ("紫金看 29-31.5 箱体，洛钼看 18-19.5 箱体", True, ACCENT_GOLD),
    ("跌到箱底缩量止跌 → 分批买    涨到箱顶放量滞涨 → 分批卖", True, WHITE),
    ("箱体不破反复做，箱体破了果断走", True, ACCENT_BLUE),
], font_size=20)

add_text_box(slide, Inches(0.8), Inches(2.6), Inches(11), Inches(0.5),
             "📝 每日操作清单", font_size=22, bold=True, color=WHITE)

checklist = [
    ("□ 9:25", "看集合竞价：高开/低开多少？有没有突破箱体？"),
    ("□ 9:30-10:30", "看前1小时量能：放量还是缩量？方向确认还是假突破？"),
    ("□ 14:30-15:00", "尾盘信号：有没有止跌形态？冲高回落还是放量突破？"),
    ("□ 收盘后", "记录今天是否触发买卖点，更新持仓计划"),
    ("□ 每周五", "复盘周线：本周是否在箱体内？箱体是否被破坏？调整下周区间"),
]
for i, (time_label, task) in enumerate(checklist):
    y = Inches(3.1) + Inches(0.42) * i
    add_text_box(slide, Inches(1.2), y, Inches(1.8), Inches(0.35), time_label, font_size=13, bold=True, color=ACCENT_GOLD)
    add_text_box(slide, Inches(3.2), y, Inches(8.5), Inches(0.35), task, font_size=13, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
             "⚡ 什么时候推翻策略重新评估？", font_size=18, bold=True, color=ACCENT_RED)
add_multiline_box(slide, Inches(1.2), Inches(5.95), Inches(10.5), Inches(1.2), [
    ("1. 收盘价有效跌破止损线（紫金27.50 / 洛钼15.80）→ 箱体已破，策略失效", False, WHITE),
    ("2. 放量突破箱体上沿（紫金34.30 / 洛钼20.00）→ 进入新趋势，需重新划定箱体", False, WHITE),
    ("3. 公司基本面出现重大变化（业绩暴雷、行业政策巨变等）", False, WHITE),
    ("4. 连续止损3次 → 说明策略在当前市场不适配，暂停并复盘", False, ACCENT_RED),
], font_size=13)

# 底部免责
add_text_box(slide, Inches(1.5), Inches(7.1), Inches(10), Inches(0.3),
             "⚠️ 以上分析基于公开技术指标和盘面数据整理，仅供思路参考，不构成投资建议。投资有风险，入市需谨慎。",
             font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── 保存 ────────────────────────────────────────────────
output_path = r"C:\Users\西西家的咩咩\箱体波动策略_紫金矿业_洛阳钼业.pptx"
prs.save(output_path)
print("PPT generated: " + output_path)
print("Total slides: " + str(len(prs.slides)))
